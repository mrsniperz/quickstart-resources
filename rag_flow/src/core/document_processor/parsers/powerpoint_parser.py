"""
模块名称: powerpoint_parser
功能描述: PowerPoint文档解析器，基于python-pptx实现PowerPoint文档的内容提取和结构分析
创建日期: 2024-01-15
作者: Sniperz
版本: v1.0.0
"""

from typing import Dict, List, Optional, Any
from pathlib import Path
from dataclasses import dataclass
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape

# 导入统一日志管理器
try:
    from src.utils.logger import SZ_LoggerManager
    logger = SZ_LoggerManager.setup_logger(__name__)
except ImportError:
    # 回退到标准logging
    import logging
    logger = logging.getLogger(__name__)

from ..extractors.metadata_extractor import MetadataExtractor


@dataclass
class PowerPointParseResult:
    """PowerPoint解析结果数据类"""
    text_content: str
    metadata: Dict[str, Any]
    slides: List[Dict[str, Any]]
    notes: List[Dict[str, Any]]
    structure_info: Dict[str, Any]


class PowerPointParser:
    """
    PowerPoint文档处理器
    
    基于python-pptx实现PowerPoint文档的全面解析，包括：
    - 幻灯片内容提取
    - 文本和形状分析
    - 备注内容提取
    - 演示文稿结构分析
    - 元数据提取
    """
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        初始化PowerPoint解析器
        
        Args:
            config (dict, optional): 配置参数
                - extract_notes (bool): 是否提取备注，默认True
                - extract_shapes (bool): 是否提取形状信息，默认True
                - preserve_slide_structure (bool): 是否保持幻灯片结构，默认True
        """
        self.config = config or {}
        self.logger = logger
        
        # 初始化提取器
        self.metadata_extractor = MetadataExtractor()
        
        # 配置参数
        self.extract_notes = self.config.get('extract_notes', True)
        self.extract_shapes = self.config.get('extract_shapes', True)
        self.preserve_slide_structure = self.config.get('preserve_slide_structure', True)
        
    def parse(self, file_path: str) -> PowerPointParseResult:
        """
        解析PowerPoint文档
        
        Args:
            file_path (str): PowerPoint文件路径
            
        Returns:
            PowerPointParseResult: 解析结果
            
        Raises:
            FileNotFoundError: 文件不存在
            ValueError: 文件格式不支持
            Exception: 解析过程中的其他错误
        """
        try:
            file_path = Path(file_path)
            if not file_path.exists():
                raise FileNotFoundError(f"PowerPoint文件不存在: {file_path}")
                
            if file_path.suffix.lower() not in ['.pptx', '.ppt']:
                raise ValueError(f"不支持的文件格式: {file_path.suffix}")
                
            self.logger.info(f"开始解析PowerPoint文档: {file_path}")
            
            # 打开演示文稿
            prs = Presentation(str(file_path))
            
            # 提取元数据
            metadata = self._extract_metadata(prs, file_path)
            
            # 提取幻灯片内容
            slides = self._extract_slides(prs)
            
            # 提取备注
            notes = []
            if self.extract_notes:
                notes = self._extract_notes(prs)
            
            # 生成文本内容
            text_content = self._generate_text_content(slides, notes)
            
            # 分析文档结构
            structure_info = self._analyze_structure(prs, slides)
            
            result = PowerPointParseResult(
                text_content=text_content,
                metadata=metadata,
                slides=slides,
                notes=notes,
                structure_info=structure_info
            )
            
            self.logger.info(f"PowerPoint解析完成: {len(slides)}张幻灯片, {len(notes)}条备注")
            return result
            
        except Exception as e:
            self.logger.error(f"PowerPoint解析失败: {e}")
            raise
    
    def _extract_metadata(self, prs: Presentation, file_path: Path) -> Dict[str, Any]:
        """
        提取演示文稿元数据
        
        Args:
            prs: python-pptx演示文稿对象
            file_path: 文件路径
            
        Returns:
            dict: 元数据信息
        """
        try:
            core_props = prs.core_properties
            
            metadata = {
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'keywords': core_props.keywords or '',
                'comments': core_props.comments or '',
                'category': core_props.category or '',
                'created': core_props.created,
                'modified': core_props.modified,
                'last_modified_by': core_props.last_modified_by or '',
                'revision': core_props.revision,
                'version': core_props.version or '',
                'file_path': str(file_path),
                'file_size': file_path.stat().st_size if file_path.exists() else 0,
                'file_extension': file_path.suffix.lower(),
                'slide_count': len(prs.slides),
                'slide_layouts_count': len(prs.slide_layouts),
                'slide_masters_count': len(prs.slide_masters),
                'is_powerpoint_document': True
            }
            
            return metadata
            
        except Exception as e:
            self.logger.error(f"元数据提取失败: {e}")
            return {'file_path': str(file_path), 'is_powerpoint_document': True}
    
    def _extract_slides(self, prs: Presentation) -> List[Dict[str, Any]]:
        """
        提取所有幻灯片内容
        
        Args:
            prs: python-pptx演示文稿对象
            
        Returns:
            list: 幻灯片信息列表
        """
        try:
            slides = []
            
            for slide_idx, slide in enumerate(prs.slides):
                try:
                    slide_info = self._extract_slide_content(slide, slide_idx)
                    slides.append(slide_info)
                    
                except Exception as e:
                    self.logger.warning(f"幻灯片{slide_idx + 1}提取失败: {e}")
                    # 添加空的幻灯片信息以保持索引一致性
                    slides.append({
                        'slide_number': slide_idx + 1,
                        'title': '',
                        'content': '',
                        'shapes': [],
                        'text_runs': [],
                        'error': str(e)
                    })
                    continue
            
            return slides
            
        except Exception as e:
            self.logger.error(f"幻灯片提取失败: {e}")
            return []
    
    def _extract_slide_content(self, slide: Slide, slide_idx: int) -> Dict[str, Any]:
        """
        提取单张幻灯片内容
        
        Args:
            slide: python-pptx幻灯片对象
            slide_idx: 幻灯片索引
            
        Returns:
            dict: 幻灯片信息
        """
        try:
            slide_info = {
                'slide_number': slide_idx + 1,
                'title': '',
                'content': '',
                'shapes': [],
                'text_runs': []
            }
            
            # 提取标题
            if slide.shapes.title:
                slide_info['title'] = slide.shapes.title.text
            
            # 提取所有文本内容
            text_parts = []
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                shape_info = {
                    'shape_type': str(shape.shape_type),
                    'text': '',
                    'paragraphs': []
                }
                
                # 提取形状中的文本
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        shape_info['paragraphs'].append(para_text)
                        text_parts.append(para_text)
                        
                        # 提取文本运行信息
                        for run in paragraph.runs:
                            if run.text.strip():
                                run_info = {
                                    'text': run.text,
                                    'font_name': run.font.name,
                                    'font_size': run.font.size.pt if run.font.size else None,
                                    'bold': run.font.bold,
                                    'italic': run.font.italic
                                }
                                slide_info['text_runs'].append(run_info)
                
                shape_info['text'] = '\n'.join(shape_info['paragraphs'])
                if shape_info['text']:
                    slide_info['shapes'].append(shape_info)
            
            slide_info['content'] = '\n'.join(text_parts)
            
            return slide_info
            
        except Exception as e:
            self.logger.error(f"幻灯片内容提取失败: {e}")
            return {
                'slide_number': slide_idx + 1,
                'title': '',
                'content': '',
                'shapes': [],
                'text_runs': [],
                'error': str(e)
            }
    
    def _extract_notes(self, prs: Presentation) -> List[Dict[str, Any]]:
        """
        提取所有幻灯片备注
        
        Args:
            prs: python-pptx演示文稿对象
            
        Returns:
            list: 备注信息列表
        """
        try:
            notes = []
            
            for slide_idx, slide in enumerate(prs.slides):
                try:
                    if slide.has_notes_slide:
                        notes_slide = slide.notes_slide
                        notes_text = notes_slide.notes_text_frame.text
                        
                        if notes_text.strip():
                            note_info = {
                                'slide_number': slide_idx + 1,
                                'notes_text': notes_text.strip()
                            }
                            notes.append(note_info)
                            
                except Exception as e:
                    self.logger.warning(f"幻灯片{slide_idx + 1}备注提取失败: {e}")
                    continue
            
            return notes
            
        except Exception as e:
            self.logger.error(f"备注提取失败: {e}")
            return []
    
    def _generate_text_content(self, slides: List[Dict[str, Any]], 
                             notes: List[Dict[str, Any]]) -> str:
        """
        生成文本内容
        
        Args:
            slides: 幻灯片信息列表
            notes: 备注信息列表
            
        Returns:
            str: 生成的文本内容
        """
        try:
            text_parts = []
            
            # 创建备注字典以便快速查找
            notes_dict = {note['slide_number']: note['notes_text'] for note in notes}
            
            for slide in slides:
                slide_text = [f"=== 幻灯片 {slide['slide_number']} ==="]
                
                # 添加标题
                if slide['title']:
                    slide_text.append(f"标题: {slide['title']}")
                
                # 添加内容
                if slide['content']:
                    slide_text.append(f"内容:\n{slide['content']}")
                
                # 添加备注
                if slide['slide_number'] in notes_dict:
                    slide_text.append(f"备注:\n{notes_dict[slide['slide_number']]}")
                
                text_parts.append('\n'.join(slide_text))
            
            return '\n\n'.join(text_parts)
            
        except Exception as e:
            self.logger.error(f"文本内容生成失败: {e}")
            return ""
    
    def _analyze_structure(self, prs: Presentation, slides: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        分析演示文稿结构
        
        Args:
            prs: python-pptx演示文稿对象
            slides: 幻灯片信息列表
            
        Returns:
            dict: 结构分析信息
        """
        try:
            structure = {
                'total_slides': len(slides),
                'slides_with_title': sum(1 for slide in slides if slide.get('title')),
                'slides_with_content': sum(1 for slide in slides if slide.get('content')),
                'total_text_runs': sum(len(slide.get('text_runs', [])) for slide in slides),
                'total_shapes': sum(len(slide.get('shapes', [])) for slide in slides),
                'slide_layouts_used': len(prs.slide_layouts),
                'slide_masters_used': len(prs.slide_masters),
                'average_content_length': 0,
                'longest_slide': None,
                'shortest_slide': None
            }
            
            # 计算内容长度统计
            content_lengths = [len(slide.get('content', '')) for slide in slides]
            if content_lengths:
                structure['average_content_length'] = sum(content_lengths) / len(content_lengths)
                
                max_length = max(content_lengths)
                min_length = min(content_lengths)
                
                structure['longest_slide'] = content_lengths.index(max_length) + 1
                structure['shortest_slide'] = content_lengths.index(min_length) + 1
            
            return structure
            
        except Exception as e:
            self.logger.error(f"结构分析失败: {e}")
            return {}
    
    def extract_slide_content(self, file_path: str, slide_number: int) -> Dict[str, Any]:
        """
        提取指定幻灯片的内容
        
        Args:
            file_path (str): PowerPoint文件路径
            slide_number (int): 幻灯片编号（从1开始）
            
        Returns:
            dict: 幻灯片内容
            
        Raises:
            IndexError: 幻灯片编号超出范围
        """
        try:
            prs = Presentation(file_path)
            
            if slide_number < 1 or slide_number > len(prs.slides):
                raise IndexError(f"幻灯片编号超出范围: {slide_number}")
            
            slide = prs.slides[slide_number - 1]
            return self._extract_slide_content(slide, slide_number - 1)
            
        except Exception as e:
            self.logger.error(f"幻灯片内容提取失败: {e}")
            raise
    
    def get_supported_formats(self) -> List[str]:
        """
        获取支持的文件格式
        
        Returns:
            list: 支持的文件扩展名列表
        """
        return ['.pptx', '.ppt']
